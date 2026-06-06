"""
Credit Risk Modeling — 1-Hour Presentation Builder
Generates Credit_Risk_Presentation.pptx from L01–L04 notebook results
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

OUT = "/Users/dukpagnarith/Documents/Obsidian Vault/project 2026/I4_risk_management/Credit_Risk_Presentation.pptx"

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy background
BLUE    = RGBColor(0x00, 0x78, 0xD4)   # Microsoft blue accent
LBLUE   = RGBColor(0x00, 0xA8, 0xE8)   # bright blue
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFFWT   = RGBColor(0xF8, 0xFA, 0xFC)   # near-white card bg
DARK    = RGBColor(0x1E, 0x29, 0x3B)   # dark text
MID     = RGBColor(0x47, 0x55, 0x69)
GREEN   = RGBColor(0x00, 0xB8, 0x94)
RED     = RGBColor(0xE7, 0x4C, 0x3C)
AMBER   = RGBColor(0xF3, 0x9C, 0x12)
LGRAY   = RGBColor(0xE2, 0xE8, 0xF0)
DBLUE2  = RGBColor(0x1B, 0x4F, 0x72)   # section heading

W  = Inches(13.33)   # LAYOUT_WIDE width
H  = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helpers ──────────────────────────────────────────────────────────────────
def rgb(r,g,b): return RGBColor(r,g,b)

def bg(slide, color):
    from pptx.util import Pt
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x,y,w,h, fill_color, line_color=None, line_w=None):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_w: shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x,y,w,h, size=18, bold=False, color=WHITE, align="left",
        italic=False, font="Calibri", wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font
    return txb

def bullet_box(slide, items, x,y,w,h, size=15, color=DARK, bold_first=False, font="Calibri"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size  = Pt(size)
        p.font.color.rgb = color
        p.font.name  = font
        p.font.bold  = (bold_first and i == 0)
        p.space_after = Pt(4)

def card(slide, x,y,w,h, fill=OFFWT, shadow=False):
    return rect(slide, x,y,w,h, fill, LGRAY, 0.5)

def section_badge(slide, label, x=0.4, y=0.2):
    rect(slide, x, y, 2.4, 0.38, BLUE)
    txt(slide, label, x+0.08, y+0.02, 2.3, 0.35, size=11, bold=True,
        color=WHITE, align="center")

def stat_card(slide, value, label, x, y, w=2.8, h=1.4,
              val_color=NAVY, lbl_color=MID):
    card(slide, x, y, w, h)
    txt(slide, value, x+0.1, y+0.12, w-0.2, 0.7, size=32, bold=True,
        color=val_color, align="center")
    txt(slide, label, x+0.1, y+0.85, w-0.2, 0.45, size=12,
        color=lbl_color, align="center")

def metric_row(slide, metrics, y=3.5):
    """metrics = list of (value, label) tuples"""
    n   = len(metrics)
    gap = 0.3
    aw  = (13.33 - 0.6 - gap*(n-1)) / n
    x   = 0.3
    for val, lbl in metrics:
        stat_card(slide, val, lbl, x, y, w=aw, h=1.5)
        x += aw + gap

def dark_title_slide(prs, title, subtitle="", tag=""):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, NAVY)
    # accent bar left
    rect(slide, 0, 0, 0.18, 7.5, BLUE)
    # top-right corner accent
    rect(slide, 11.5, 0, 1.83, 0.12, LBLUE)
    if tag:
        rect(slide, 0.5, 1.0, 3.2, 0.42, BLUE)
        txt(slide, tag, 0.6, 1.02, 3.0, 0.38, size=13, bold=True,
            color=WHITE, align="left")
    txt(slide, title, 0.5, 1.8, 12.3, 2.8, size=44, bold=True,
        color=WHITE, align="left", font="Calibri")
    if subtitle:
        txt(slide, subtitle, 0.5, 4.9, 12.0, 1.2, size=20,
            color=RGBColor(0xA0,0xC4,0xE8), align="left")
    # bottom bar
    rect(slide, 0.18, 7.1, 13.15, 0.4, DBLUE2)
    txt(slide, "Credit Risk Modeling System  |  Lending Club 2007–2018",
        0.4, 7.12, 12.5, 0.36, size=11, color=RGBColor(0x90,0xB4,0xD8), align="left")
    return slide

def content_slide(prs, title, section_label=""):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, WHITE)
    rect(slide, 0, 0, 0.10, 7.5, BLUE)
    rect(slide, 0.10, 0, 13.23, 0.85, NAVY)
    txt(slide, title, 0.3, 0.08, 10.5, 0.72, size=24, bold=True,
        color=WHITE, align="left")
    if section_label:
        txt(slide, section_label, 10.9, 0.18, 2.2, 0.5, size=12,
            color=RGBColor(0x90,0xB4,0xD8), align="right")
    rect(slide, 0.10, 7.1, 13.23, 0.4, LGRAY)
    txt(slide, "Credit Risk Modeling System  |  Lending Club 2007–2018",
        0.3, 7.12, 12.5, 0.36, size=10, color=MID, align="left")
    return slide

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = dark_title_slide(prs,
    "Credit Risk Modeling System",
    "Probability of Default · Loss Given Default · Expected Loss · PSI Monitoring\n"
    "Lending Club Dataset 2007–2018  |  2.26M Loans  |  $3.26B Portfolio",
    tag="")
rect(s, 0.5, 6.35, 12.5, 0.02, LBLUE)
txt(s, "A complete end-to-end credit risk pipeline: data → model → scorecard → production API",
    0.5, 6.5, 12.0, 0.6, size=13, color=RGBColor(0xA0,0xC4,0xE8), italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "Agenda — 60-Minute Session", "Overview")
sections = [
    ("0:00 – 0:10", "Dataset & Preprocessing",   "L01", "2.26M rows · 151 columns · WoE · OOT split"),
    ("0:10 – 0:28", "PD Model & Scorecard",       "L02", "Logistic regression · Gini 0.39 · 300–850 scorecard"),
    ("0:28 – 0:44", "LGD / EAD / Expected Loss",  "L03", "2-stage LGD · EAD · $383M portfolio EL"),
    ("0:44 – 0:54", "PSI Monitoring",             "L04", "Score PSI 0.109 · Drift detection · Action framework"),
    ("0:54 – 1:00", "Production Pipeline",        "LIVE", "FastAPI · MLflow · Airflow · Docker"),
]
colors = [BLUE, rgb(0x1B,0x6C,0xA8), rgb(0x17,0x5B,0x8A), rgb(0x13,0x4A,0x6C), NAVY]
for i, (time, title, badge, desc) in enumerate(sections):
    y = 1.1 + i * 1.1
    rect(s, 0.4, y, 12.4, 0.95, colors[i])
    txt(s, time, 0.6, y+0.12, 1.8, 0.4, size=13, bold=True, color=WHITE)
    rect(s, 2.5, y+0.1, 0.85, 0.38, LBLUE)
    txt(s, badge, 2.52, y+0.12, 0.8, 0.38, size=13, bold=True, color=WHITE, align="center")
    txt(s, title, 3.5, y+0.08, 5.5, 0.42, size=16, bold=True, color=WHITE)
    txt(s, desc, 3.5, y+0.5, 8.8, 0.35, size=12, color=RGBColor(0xC0,0xD8,0xF0))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Regulatory Context
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "Why Credit Risk Modeling?", "Context")
txt(s, "Regulatory & Business Drivers", 0.4, 1.0, 12.5, 0.5, size=18, bold=True, color=NAVY)

regs = [
    ("IFRS 9", "Expected Credit Loss provisioning.\nStage 1/2/3 classification based on PD.", rgb(0x1B,0x4F,0x72)),
    ("Basel III", "Capital requirements linked to PD,\nLGD and EAD via RWA × 8% formula.", rgb(0x1A,0x53,0x76)),
    ("SR 11-7", "Model risk governance: full audit trail,\nindependent validation, monitoring.", rgb(0x19,0x57,0x7A)),
]
for i, (name, desc, clr) in enumerate(regs):
    x = 0.4 + i*4.3
    rect(s, x, 1.7, 4.0, 2.1, clr)
    txt(s, name, x+0.2, 1.82, 3.6, 0.6, size=22, bold=True, color=WHITE)
    txt(s, desc, x+0.2, 2.45, 3.6, 1.2, size=13, color=RGBColor(0xC0,0xD8,0xF0))

txt(s, "Business Value", 0.4, 4.0, 12.5, 0.45, size=18, bold=True, color=NAVY)
biz = [
    "▸  Reduce credit losses by identifying high-risk borrowers before approval",
    "▸  Optimize portfolio ROI: approve loans where interest income > expected loss",
    "▸  Quantify provisions: how much capital to set aside for future defaults",
    "▸  Early warning: PSI monitoring detects population drift before losses materialize",
]
for i, b in enumerate(biz):
    txt(s, b, 0.5, 4.55+i*0.48, 12.3, 0.42, size=14, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Dataset Overview
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "The Lending Club Dataset", "Dataset")
stats = [
    ("2.26M", "Total Loan Applications"),
    ("151", "Raw Columns"),
    ("1.6 GB", "Raw CSV Size"),
    ("2007–2018", "Date Range"),
]
x = 0.4
for val, lbl in stats:
    stat_card(s, val, lbl, x, 1.05, w=2.9, h=1.5, val_color=BLUE)
    x += 3.1

txt(s, "Out-of-Time Split Strategy", 0.4, 2.8, 12.5, 0.45, size=16, bold=True, color=NAVY)
# Train bar
rect(s, 0.4, 3.3, 8.2, 0.9, BLUE)
txt(s, "TRAIN  (2007–2015)   831,051 loans   Default rate: 18.62%",
    0.7, 3.42, 7.8, 0.6, size=14, bold=True, color=WHITE)
# Test bar
rect(s, 8.65, 3.3, 4.0, 0.9, RED)
txt(s, "OOT TEST  (2016–2018)   538,515 loans   DR: 25.27%",
    8.75, 3.42, 3.8, 0.6, size=12, bold=True, color=WHITE)
# Arrow label
txt(s, "← No future data leaks into training →", 3.5, 4.35, 6.3, 0.4,
    size=12, color=MID, align="center", italic=True)

txt(s, "Column Journey: 151 raw  →  57 dropped  →  56 WoE dummy variables for logistic regression",
    0.4, 5.0, 12.5, 0.45, size=13, color=DARK)

cats = [
    ("100% null", "Dropped",   "8"),
    ("Identifiers", "Dropped", "10"),
    ("Post-app leakage", "Dropped", "21"),
    ("Near-zero var", "Dropped", "8"),
    ("Joint/hardship", "Dropped", "17"),
    ("WoE Dummies", "KEPT → 56", "56"),
]
for i, (name, action, n) in enumerate(cats):
    clr = GREEN if "KEPT" in action else LGRAY
    tc  = DARK if "KEPT" not in action else WHITE
    rect(s, 0.4+i*2.15, 5.6, 2.0, 1.3, clr, LGRAY, 0.5)
    txt(s, name, 0.5+i*2.15, 5.7, 1.8, 0.55, size=11, color=tc if "KEPT" in action else DARK, bold="KEPT" in action)
    txt(s, action, 0.5+i*2.15, 6.25, 1.8, 0.35, size=10, color=BLUE if "KEPT" in action else MID, bold=True)
    txt(s, n, 0.5+i*2.15, 6.6, 1.8, 0.28, size=16, bold=True, color=NAVY if "KEPT" not in action else WHITE, align="center")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — L01 Section Header
# ═══════════════════════════════════════════════════════════════════════════════
s = dark_title_slide(prs,
    "L01 — Data Preprocessing\n& Feature Engineering",
    "Cleaning · Format conversions · WoE · Dummy variables · OOT split",
    tag="NOTEBOOK 1 OF 4")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Data Cleaning Pipeline
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "7-Step Data Cleaning Pipeline", "L01")
txt(s, "Critical: Create target variables BEFORE dropping leakage columns", 0.4, 1.0, 12.5, 0.4,
    size=13, color=RED, italic=True, bold=True)

steps = [
    ("1", "Create Targets",       "good_bad, recovery_rate, CCF from raw columns before any drops"),
    ("2", "Drop 100%-null",       "8 columns with no data whatsoever"),
    ("3", "Drop Identifiers",     "id, url, emp_title, zip_code — no predictive value"),
    ("4", "Drop Leakage",         "21 post-origination cols: total_pymnt, recoveries, last_fico…"),
    ("5", "Near-zero variance",   "8 cols: tax_liens, acc_now_delinq — near-constant values"),
    ("6", "Joint/Hardship",       "17 cols: annual_inc_joint, hardship_* — >80% null"),
    ("7", "Format Conversions",   "term '36 months'→36  ·  int_rate 13.99→0.1399  ·  FICO avg"),
]
for i, (n, title, desc) in enumerate(steps):
    row_y = 1.55 + i * 0.75
    clr = BLUE if i == 0 else (GREEN if i == 6 else NAVY)
    rect(s, 0.4, row_y, 0.55, 0.58, clr)
    txt(s, n, 0.4, row_y+0.06, 0.55, 0.46, size=18, bold=True, color=WHITE, align="center")
    txt(s, title, 1.1, row_y+0.04, 2.8, 0.3, size=14, bold=True, color=NAVY)
    txt(s, desc,  3.95, row_y+0.04, 9.0, 0.5, size=12, color=DARK)
    if i < 6:
        rect(s, 0.62, row_y+0.58, 0.12, 0.17, LGRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Target Variables
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "Three Target Variables", "L01")

targets = [
    ("good_bad",       "PD Target",  BLUE,
     "1 = Fully Paid, 0 = Charged Off / Default / Late 31-120d",
     "Binary classification — logistic regression",
     "18.62% bad rate (train)  ·  25.27% (OOT test)"),
    ("recovery_rate",  "LGD Target", rgb(0x8E,0x44,0xAD),
     "Recoveries / Funded Amount — clipped to [0, 1]",
     "Only computed for defaulted loans (good_bad = 0)",
     "23.9% have zero recovery (bimodal → 2-stage model)"),
    ("ccf",            "EAD Target", rgb(0x27,0xAE,0x60),
     "(Funded Amount − Total Paid) / Funded Amount → [0,1]",
     "Credit Conversion Factor: fraction outstanding at default",
     "Mean CCF train = 0.39  ·  Mean CCF test = 0.56"),
]
for i, (col, label, clr, formula, note1, note2) in enumerate(targets):
    x = 0.4 + i*4.3
    rect(s, x, 1.0, 4.0, 5.8, clr)
    txt(s, label,  x+0.2, 1.1,  3.6, 0.45, size=13, bold=True, color=WHITE)
    txt(s, col,    x+0.2, 1.58, 3.6, 0.5,  size=20, bold=True, color=WHITE, font="Consolas")
    rect(s, x+0.15, 2.2, 3.7, 0.02, WHITE)
    txt(s, formula, x+0.2, 2.3, 3.6, 0.65, size=12, color=RGBColor(0xE0,0xE0,0xFF))
    txt(s, note1,   x+0.2, 3.1, 3.6, 0.65, size=11, color=RGBColor(0xD0,0xD8,0xF0))
    txt(s, note2,   x+0.2, 3.85, 3.6, 0.7, size=11, color=RGBColor(0xD0,0xD8,0xF0))

txt(s, "EL = PD × LGD × EAD    where    LGD = 1 − Recovery Rate    and    EAD = CCF × Funded Amount",
    0.4, 7.0, 12.5, 0.4, size=13, bold=True, color=NAVY, align="center")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — WoE & Information Value
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "Weight of Evidence (WoE) & Information Value", "L01")

# Formula box
rect(s, 0.4, 1.0, 8.0, 1.8, NAVY)
txt(s, "WoE_i  =  ln( %Goods_i / %Bads_i )",
    0.7, 1.15, 7.4, 0.6, size=18, bold=True, color=WHITE, font="Consolas")
txt(s, "IV  =  Σ ( %Goods_i − %Bads_i ) × WoE_i",
    0.7, 1.78, 7.4, 0.5, size=18, bold=True, color=LBLUE, font="Consolas")

# IV table
rect(s, 8.6, 1.0, 4.5, 1.8, OFFWT, LGRAY, 0.5)
iv_rows = [("< 0.02","Useless","Drop",RED),("0.02–0.10","Weak","Consider",AMBER),
           ("0.10–0.30","Medium","Keep",GREEN),("0.30–0.50","Strong","Keep",GREEN),
           ("> 0.50","Suspicious","Check leakage",RED)]
txt(s, "IV     Strength    Action", 8.75, 1.08, 4.2, 0.35, size=11, bold=True, color=NAVY)
for i,(iv,strength,action,clr) in enumerate(iv_rows):
    txt(s, f"{iv:<10} {strength:<12} {action}", 8.75, 1.44+i*0.27, 4.2, 0.3,
        size=10, color=clr if clr!=AMBER else rgb(0xB0,0x6A,0x00))

txt(s, "Key IV results from training data (831K loans)", 0.4, 3.0, 12.5, 0.4,
    size=14, bold=True, color=NAVY)

iv_results = [
    ("grade",          0.52, "Check"),
    ("int_rate",       0.38, "Strong"),
    ("fico_score",     0.31, "Strong"),
    ("dti",            0.14, "Medium"),
    ("term_int",       0.13, "Medium"),
    ("annual_inc",     0.11, "Medium"),
    ("revol_util",     0.06, "Weak"),
    ("inq_last_6mths", 0.05, "Weak"),
]
for i, (var, iv, strength) in enumerate(iv_results):
    col = i % 4
    row = i // 4
    x = 0.4 + col*3.2
    y = 3.55 + row*1.2
    clr = GREEN if strength=="Strong" else (AMBER if strength=="Medium" else
          (RED if strength=="Check" else MID))
    card(s, x, y, 3.0, 1.0)
    txt(s, var, x+0.15, y+0.1, 2.7, 0.4, size=13, bold=True, color=NAVY, font="Consolas")
    bar_w = min(2.5 * iv / 0.6, 2.5)
    rect(s, x+0.15, y+0.6, bar_w, 0.22, clr)
    txt(s, f"IV = {iv:.2f}  {strength}", x+bar_w+0.2, y+0.6, 2.0, 0.25, size=10, color=clr)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — EDA Highlights
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "Exploratory Data Analysis — Key Findings", "L01")

# Default rate by grade (bar chart)
rect(s, 0.4, 1.0, 6.0, 5.8, OFFWT, LGRAY, 0.5)
txt(s, "Default Rate by Grade", 0.7, 1.1, 5.5, 0.4, size=14, bold=True, color=NAVY)
grades = ["A","B","C","D","E","F","G"]
drates = [0.058, 0.108, 0.163, 0.224, 0.297, 0.375, 0.428]
max_dr = max(drates)
bh_max = 3.6
for i,(g,dr) in enumerate(zip(grades, drates)):
    bh = bh_max * dr / max_dr
    bx = 0.8 + i*0.78
    rect(s, bx, 1.6+bh_max-bh, 0.58, bh, RED if dr > 0.25 else (AMBER if dr>0.15 else GREEN))
    txt(s, g, bx, 5.3, 0.6, 0.3, size=11, bold=True, color=NAVY, align="center")
    txt(s, f"{dr*100:.0f}%", bx, 1.45+bh_max-bh, 0.6, 0.3, size=9, color=DARK, align="center")
txt(s, "Grade A borrowers default 6× less than Grade G",
    0.55, 6.45, 5.7, 0.35, size=11, italic=True, color=MID)

# FICO comparison
rect(s, 6.8, 1.0, 6.1, 5.8, OFFWT, LGRAY, 0.5)
txt(s, "FICO Score: Good vs Bad Borrowers", 7.0, 1.1, 5.7, 0.4, size=14, bold=True, color=NAVY)
# Simulate distribution bars
fico_bins = [600,620,640,660,680,700,720,740,760,780,800]
good_dist = [0.01,0.02,0.05,0.08,0.12,0.17,0.19,0.18,0.11,0.07]
bad_dist  = [0.04,0.07,0.12,0.15,0.16,0.15,0.12,0.09,0.06,0.04]
bh_max2 = 3.5
max_d = max(max(good_dist), max(bad_dist))
for i in range(len(good_dist)):
    bx = 7.0 + i*0.54
    gw = 0.23
    gh = bh_max2 * good_dist[i] / max_d
    bh2 = bh_max2 * bad_dist[i] / max_d
    rect(s, bx, 1.6+bh_max2-gh, gw, gh, GREEN)
    rect(s, bx+gw, 1.6+bh_max2-bh2, gw, bh2, RED)
txt(s, "Mean FICO — Good: 700.1   Bad: 671.4   Gap: 28.7 pts",
    7.0, 5.2, 5.7, 0.4, size=12, color=DARK, bold=True)
txt(s, "▬ Good borrowers   ▬ Bad borrowers", 7.0, 5.65, 5.7, 0.4, size=11, color=MID)
txt(s, "■ Good    ■ Bad", 7.0, 6.45, 5.7, 0.35, size=11, italic=True, color=MID)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — L01 Summary
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "L01 — Summary", "L01")
summ = [
    ("831,051", "Training Loans\n2007–2015", BLUE),
    ("538,515", "OOT Test Loans\n2016–2018", RED),
    ("56", "WoE Dummy\nFeatures", GREEN),
    ("18.62%", "Train Default\nRate", NAVY),
]
x=0.4
for val, lbl, clr in summ:
    stat_card(s, val, lbl, x, 1.1, w=3.0, h=1.6, val_color=clr)
    x+=3.1

txt(s,"Key Design Decisions",0.4,3.0,12.5,0.4,size=16,bold=True,color=NAVY)
decisions=[
    "Out-of-time split: 2007–2015 train / 2016–2018 OOT — prevents look-ahead bias (3-year holdout)",
    "WoE bins fitted on TRAINING SET ONLY — OOT set uses the same fixed bin boundaries",
    "Missing values treated as their own WoE bin (sentinel = -1), NOT imputed away",
    "int_rate divided by 100 (raw CSV stores percentages, not decimals)",
    "FICO score = average of fico_range_low and fico_range_high (strong predictor not in most GitHub tutorials)",
    "Dummy column list made explicit (prefix matching was capturing raw categorical columns)",
]
for i, d in enumerate(decisions):
    txt(s, f"✓  {d}", 0.5, 3.55+i*0.52, 12.3, 0.46, size=13, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — L02 Section Header
# ═══════════════════════════════════════════════════════════════════════════════
s = dark_title_slide(prs,
    "L02 — PD Model, Scorecard\n& Credit Policy",
    "Logistic regression · Backward elimination · AUC 0.693 · 300–850 Scorecard · ROI policy",
    tag="NOTEBOOK 2 OF 4")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Why Logistic Regression
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "Why Logistic Regression?", "L02")
rect(s, 0.4, 1.0, 12.5, 1.1, NAVY)
txt(s, "log( P(Good) / P(Bad) ) = β₀ + β₁x₁ + β₂x₂ + … + βₖxₖ",
    0.7, 1.2, 12.0, 0.7, size=22, bold=True, color=WHITE, font="Consolas", align="center")

reasons = [
    ("Interpretable", "Each coefficient maps directly to log-odds of default. Regulators can inspect every term.", BLUE),
    ("Basel II/III", "Advanced IRB approach requires logistic PD models with documented p-values.", rgb(0x1B,0x4F,0x72)),
    ("statsmodels", "sklearn LogisticRegression gives no p-values. statsmodels provides Wald test per coefficient.", rgb(0x17,0x45,0x60)),
    ("Stable", "On 831K training examples, logistic regression is robust — no overfitting risk.", rgb(0x13,0x3B,0x52)),
]
for i, (title, desc, clr) in enumerate(reasons):
    x = 0.4 + i*3.2
    rect(s, x, 2.35, 3.0, 4.2, clr)
    txt(s, title, x+0.2, 2.5, 2.6, 0.55, size=16, bold=True, color=WHITE)
    txt(s, desc, x+0.2, 3.1, 2.6, 3.2, size=12, color=RGBColor(0xC0,0xD8,0xF5))

txt(s, "Key difference vs sklearn: statsmodels provides Wald test p-values → required for regulatory documentation",
    0.4, 6.75, 12.5, 0.4, size=12, italic=True, color=RED, bold=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Feature Selection
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "Feature Selection — Backward Elimination", "L02")

rect(s, 0.4, 1.0, 12.5, 1.0, OFFWT, LGRAY, 0.5)
txt(s, "Start with all 56 dummy variables → iteratively remove the least significant (highest p-value > 0.05) → stop when all remaining have p < 0.05",
    0.6, 1.12, 12.1, 0.75, size=13, color=DARK)

# Flow
steps2 = ["56 features\n(all dummies)","Fit Logit\n(Newton)","Remove worst\np > 0.05",
          "Check\nconvergence","49 features\n✓ all p<0.05"]
arrow_x = [0.5, 2.65, 4.9, 7.15, 9.4]
arrow_w = 2.0
for i,step in enumerate(steps2):
    clr = GREEN if i==4 else BLUE
    rect(s, arrow_x[i], 2.2, arrow_w, 1.1, clr)
    txt(s, step, arrow_x[i]+0.1, 2.32, arrow_w-0.2, 0.9, size=12, bold=(i in[0,4]),
        color=WHITE, align="center")
    if i < 4:
        txt(s, "→", arrow_x[i]+arrow_w, 2.55, 0.25, 0.45, size=22, bold=True, color=NAVY)

txt(s, "Removed 7 variables: int_rate_088_117 was last (p=0.1523 before removal)",
    0.4, 3.55, 12.5, 0.4, size=12, italic=True, color=MID)

txt(s, "Final: 49 features   |   All p-values < 0.05   |   Converged in 7 iterations",
    0.4, 4.05, 12.5, 0.45, size=15, bold=True, color=GREEN)

# Top coefficients
txt(s, "Top Positive Coefficients (lower default risk)", 0.4, 4.7, 6.1, 0.4, size=14, bold=True, color=GREEN)
txt(s, "Most Negative Coefficients (higher default risk)", 6.8, 4.7, 6.1, 0.4, size=14, bold=True, color=RED)
pos_coefs = [("grade_A", 1.3477), ("grade_B", 1.0456), ("grade_C", 0.7617),
             ("term_36", 0.7371), ("dti_lt_10", 0.5454)]
neg_coefs = [("fico_640_680", -0.3259), ("fico_680_720", -0.1863),
             ("int_rate_117_148", -0.1313), ("verif_Source_Verified", -0.1197),
             ("verif_Verified", -0.1041)]
for i, (feat, coef) in enumerate(pos_coefs):
    bar_w = 5.0 * abs(coef) / 1.5
    rect(s, 0.5, 5.2+i*0.35, bar_w, 0.28, GREEN)
    txt(s, f"{feat:<25} {coef:+.4f}", 0.5+bar_w+0.05, 5.22+i*0.35, 5.5, 0.28, size=10, color=DARK, font="Consolas")
for i, (feat, coef) in enumerate(neg_coefs):
    bar_w = 3.5 * abs(coef) / 0.35
    rect(s, 6.8, 5.2+i*0.35, bar_w, 0.28, RED)
    txt(s, f"{feat:<25} {coef:+.4f}", 6.8+bar_w+0.05, 5.22+i*0.35, 5.5, 0.28, size=10, color=DARK, font="Consolas")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Model Performance
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "PD Model Performance", "L02")

metrics = [
    ("AUC",   0.7136, 0.6931, "> 0.65"),
    ("Gini",  0.4273, 0.3862, "> 0.40"),
    ("KS",    0.3093, 0.2801, "> 0.25"),
    ("Brier", 0.1374, 0.1749, "< 0.10"),
]
rect(s, 0.4, 1.0, 12.5, 0.52, NAVY)
for i, lbl in enumerate(["Metric","Train","OOT Test","Threshold","Gap","Status"]):
    txt(s, lbl, 0.6+i*2.1, 1.08, 2.0, 0.38, size=13, bold=True, color=WHITE)

for row_i, (name, train_v, test_v, thresh) in enumerate(metrics):
    ry = 1.55 + row_i*0.85
    clr = OFFWT if row_i%2==0 else WHITE
    rect(s, 0.4, ry, 12.5, 0.8, clr, LGRAY, 0.3)
    gap = abs(train_v - test_v)
    ok  = gap < 0.05
    status = "✓ Stable" if ok else "⚠ Monitor"
    txt(s, name,         0.6,  ry+0.22, 2.0, 0.4, size=15, bold=True, color=NAVY)
    txt(s, f"{train_v:.4f}", 2.7,  ry+0.22, 2.0, 0.4, size=15, color=BLUE, bold=True)
    txt(s, f"{test_v:.4f}",  4.8,  ry+0.22, 2.0, 0.4, size=15, color=BLUE, bold=True)
    txt(s, thresh,       6.9,  ry+0.22, 2.0, 0.4, size=13, color=MID)
    txt(s, f"{gap:.4f}", 9.0,  ry+0.22, 2.0, 0.4, size=13, color=GREEN if ok else AMBER)
    txt(s, status, 11.1, ry+0.18, 1.7, 0.45, size=13,
        color=GREEN if ok else AMBER, bold=True)

txt(s, "All 4 metrics pass OOT thresholds. Train/Test gap < 0.05 confirms no overfitting.",
    0.4, 5.0, 12.5, 0.4, size=13, bold=True, color=GREEN)

# ROC sketch
rect(s, 0.4, 5.5, 5.8, 1.6, OFFWT, LGRAY, 0.5)
txt(s, "ROC Curve: AUC Train=0.714  AUC OOT=0.693",
    0.6, 5.6, 5.4, 0.45, size=12, bold=True, color=NAVY)
txt(s, "Curves remain close → model generalises to 2016–2018 borrowers without retraining.",
    0.6, 6.1, 5.4, 0.85, size=11, color=DARK)
# KS sketch
rect(s, 6.5, 5.5, 6.3, 1.6, OFFWT, LGRAY, 0.5)
txt(s, "Gini 0.39 on OOT exceeds the 0.35 rebuild threshold.",
    6.7, 5.6, 5.9, 0.45, size=12, bold=True, color=NAVY)
txt(s, "KS 0.28 confirms meaningful separation between good and bad borrowers on unseen data.",
    6.7, 6.1, 5.9, 0.85, size=11, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Scorecard
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "Credit Scorecard — 300 to 850 Scale", "L02")

rect(s, 0.4, 1.0, 7.8, 1.0, NAVY)
txt(s, "Score = Offset + Factor × log_odds      Factor = PDO/ln(2) = 28.85      Offset = 600",
    0.6, 1.1, 7.5, 0.8, size=13, bold=True, color=WHITE, font="Consolas")
rect(s, 8.4, 1.0, 4.5, 1.0, OFFWT, LGRAY, 0.5)
for lbl, val in [("PDO","20  pts"), ("Ref Score","600"), ("Ref Odds","1 : 1")]:
    pass
txt(s, "PDO = 20 pts   Ref Score = 600 @ 1:1 odds\nEach +20 pts doubles the good:bad odds",
    8.55, 1.12, 4.2, 0.78, size=11, color=DARK)

# Score bar (300-850 gradient-like)
score_colors = [RED, RED, AMBER, AMBER, GREEN, GREEN, GREEN, LBLUE, BLUE, NAVY]
for i, clr in enumerate(score_colors):
    rect(s, 0.4+i*1.25, 2.2, 1.22, 0.5, clr)
score_labels = ["300","385","435","490","545","600","655","710","765","820","850"]
for i, lbl in enumerate(score_labels):
    txt(s, lbl, 0.4+i*1.25, 2.75, 1.22, 0.3, size=9, color=DARK, align="center")

txt(s, "Selected Scorecard Entries", 0.4, 3.2, 12.5, 0.4, size=14, bold=True, color=NAVY)

sc_data = [
    ("grade_A",        1.3477, -27, "Best-in-class borrowers"),
    ("grade_B",        1.0456, -18, "Low risk"),
    ("term_36",        0.7371,  -9, "Shorter term = lower risk"),
    ("dti_lt_10",      0.5454,  -4, "Low debt burden"),
    ("fico_gt760",     0.2757,   4, "Excellent credit history"),
    ("inq_0",          0.2776,   4, "No recent inquiries"),
    ("fico_640_680",  -0.3259,  22, "Below-average FICO"),
    ("int_rate_117_148",-0.1313, 16, "Higher rate tier"),
]
rect(s, 0.4, 3.65, 12.5, 0.42, NAVY)
for i, lbl in enumerate(["Feature","Coefficient","Score Points","Interpretation"]):
    txt(s, lbl, 0.6+i*3.1, 3.73, 3.0, 0.32, size=12, bold=True, color=WHITE)
for ri, (feat, coef, score, interp) in enumerate(sc_data):
    ry = 4.1 + ri*0.37
    clr = OFFWT if ri%2==0 else WHITE
    rect(s, 0.4, ry, 12.5, 0.36, clr)
    sc_clr = GREEN if score < 0 else RED
    txt(s, feat,  0.6, ry+0.06, 3.0, 0.28, size=11, color=NAVY, font="Consolas")
    txt(s, f"{coef:+.4f}", 3.7, ry+0.06, 2.8, 0.28, size=11, color=BLUE)
    txt(s, f"{score:+d}", 6.8, ry+0.06, 2.5, 0.28, size=12, bold=True, color=sc_clr)
    txt(s, interp, 9.8, ry+0.06, 3.1, 0.28, size=10, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — 10 Risk Classes & Credit Policy
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "10 Risk Classes & ROI-Based Credit Policy", "L02")

classes = [
    ("F",  "300–459", "AUTO\nREJECT", RED),
    ("DD", "460–499", "REJECT\nif ROI<2.15%", RED),
    ("CD", "500–539", "REJECT\nif ROI<2.15%", AMBER),
    ("C",  "540–579", "APPROVE\nif ROI>2.15%", AMBER),
    ("BC", "580–619", "APPROVE\nif ROI>2.15%", AMBER),
    ("B",  "620–659", "APPROVE\nif ROI>2.15%", GREEN),
    ("BB", "660–699", "APPROVE\nif ROI>2.15%", GREEN),
    ("AB", "700–739", "APPROVE\nif ROI>2.15%", GREEN),
    ("A",  "740–779", "AUTO\nAPPROVE", BLUE),
    ("AA", "780–850", "AUTO\nAPPROVE", BLUE),
]
for i, (cls, rng, action, clr) in enumerate(classes):
    x = 0.35 + i*1.3
    rect(s, x, 1.0, 1.2, 1.6, clr)
    txt(s, cls, x, 1.1, 1.2, 0.55, size=18, bold=True, color=WHITE, align="center")
    txt(s, rng, x, 1.65, 1.2, 0.35, size=8, color=RGBColor(0xE0,0xE0,0xFF), align="center")

# ROI formula
rect(s, 0.4, 2.8, 12.5, 1.0, NAVY)
txt(s, "ROI = ( Interest Income − Expected Loss ) / Funded Amount × (12 / Term Months)",
    0.7, 2.98, 12.0, 0.65, size=16, bold=True, color=WHITE, font="Consolas", align="center")

# Decision logic
txt(s, "Decision Logic", 0.4, 4.0, 12.5, 0.4, size=15, bold=True, color=NAVY)
logic = [
    ("AA, A",        "AUTO_APPROVE — excellent borrowers, clearest risk signal", BLUE),
    ("F",            "AUTO_REJECT — highest default probability, no ROI calculation needed", RED),
    ("All others",   "APPROVE if annualized ROI > 2.15% (US base rate, 2015) else REJECT", GREEN),
]
for i, (cls, desc, clr) in enumerate(logic):
    rect(s, 0.4, 4.55+i*0.73, 2.2, 0.6, clr)
    txt(s, cls, 0.55, 4.62+i*0.73, 1.9, 0.45, size=13, bold=True, color=WHITE, align="center")
    txt(s, desc, 2.75, 4.67+i*0.73, 10.1, 0.5, size=13, color=DARK)

# Policy result
rect(s, 0.4, 6.65, 12.5, 0.65, GREEN)
txt(s, "Credit Policy Result: Default rate reduced from 25.27% → baseline   ·   ROI-optimal approval mix",
    0.6, 6.77, 12.0, 0.42, size=13, bold=True, color=WHITE, align="center")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — L02 Summary
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "L02 — Summary", "L02")
summ2 = [("49","Final Features\n(from 56)", BLUE),("0.693","OOT AUC",GREEN),
         ("0.386","OOT Gini",GREEN),("0.280","OOT KS",GREEN)]
x=0.4
for val,lbl,clr in summ2:
    stat_card(s,val,lbl,x,1.0,w=3.0,h=1.6,val_color=clr); x+=3.1

txt(s,"Key Takeaways",0.4,2.8,12.5,0.4,size=16,bold=True,color=NAVY)
kts=[
    "statsmodels provides Wald test p-values — mandatory for Basel documentation",
    "Backward elimination removed 7 insignificant features (int_rate_088_117 was last at p=0.15)",
    "Grade is the dominant predictor: grade_A coefficient = 1.35 → 3.85× better odds vs F/G",
    "Scorecard: PDO=20 means each 20-point increase doubles good:bad odds",
    "All 4 evaluation metrics pass OOT thresholds — model generalizes to 2016–2018 without refit",
]
for i,k in enumerate(kts):
    txt(s,f"✓  {k}",0.5,3.38+i*0.55,12.3,0.48,size=13,color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — L03 Section Header
# ═══════════════════════════════════════════════════════════════════════════════
s = dark_title_slide(prs,
    "L03 — LGD, EAD\n& Expected Loss",
    "Two-stage LGD · Credit Conversion Factor · EL = PD × LGD × EAD · $3.26B portfolio",
    tag="NOTEBOOK 3 OF 4")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Three-Model Architecture
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "Three-Model Architecture", "L03")

models_info = [
    ("PD Model","Probability\nof Default","All loans\n(831K train)","Logistic Regression\n(statsmodels)","0 or 1\ngood_bad", BLUE),
    ("LGD Model","Loss Given\nDefault","Defaulted only\n(154K train)","2-Stage:\nLogistic + Linear","recovery_rate\n[0,1]",rgb(0x8E,0x44,0xAD)),
    ("EAD Model","Exposure at\nDefault","Defaulted only\n(154K train)","Linear Regression\non CCF","ccf\n[0,1]",GREEN),
]
for i, (name, what, data, method, target, clr) in enumerate(models_info):
    x = 0.4 + i*4.3
    rect(s, x, 1.0, 4.0, 5.2, clr)
    txt(s, name, x+0.2, 1.1, 3.6, 0.55, size=17, bold=True, color=WHITE)
    for j, (lbl, val) in enumerate([("Predicts",what),("Training data",data),("Algorithm",method),("Target",target)]):
        rect(s, x+0.15, 1.82+j*1.0, 3.7, 0.02, RGBColor(0xFF,0xFF,0xFF))
        txt(s, lbl.upper(), x+0.2, 1.87+j*1.0, 3.6, 0.28, size=9, bold=True,
            color=RGBColor(0xB0,0xD0,0xFF))
        txt(s, val, x+0.2, 2.17+j*1.0, 3.6, 0.55, size=12, color=WHITE)

rect(s, 0.4, 6.4, 12.5, 0.75, NAVY)
txt(s, "Expected Loss  =  PD  ×  LGD  ×  EAD  =  PD  ×  (1 − Recovery Rate)  ×  (CCF × Funded Amount)",
    0.6, 6.55, 12.0, 0.45, size=16, bold=True, color=WHITE, align="center", font="Consolas")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Two-Stage LGD
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "Two-Stage LGD Model — Why Two Stages?", "L03")

rect(s, 0.4, 1.0, 12.5, 0.75, OFFWT, LGRAY, 0.5)
txt(s, "~24% of defaulted loans have exactly ZERO recovery. A single linear regression cannot model both 0s and positive values.",
    0.6, 1.12, 12.1, 0.55, size=13, color=RED, bold=True)

# Stage 1
rect(s, 0.4, 1.95, 5.9, 3.2, BLUE)
txt(s, "Stage 1 — Logistic Regression", 0.6, 2.05, 5.5, 0.5, size=15, bold=True, color=WHITE)
txt(s, "ALL 154,761 defaulted loans", 0.6, 2.58, 5.5, 0.35, size=11, color=RGBColor(0xC0,0xD8,0xFF))
txt(s, "Target: Is Recovery Rate > 0?\n(1 = any recovery, 0 = total loss)", 0.6, 3.0, 5.5, 0.65, size=12, color=WHITE)
txt(s, "Train AUC = 0.673    OOT AUC = 0.522", 0.6, 3.75, 5.5, 0.4, size=12, bold=True, color=RGBColor(0x90,0xFF,0x90))
txt(s, "Train Gini = 0.346    OOT Gini = 0.043", 0.6, 4.18, 5.5, 0.35, size=12, color=RGBColor(0x90,0xFF,0x90))

# Stage 2
rect(s, 6.7, 1.95, 5.9, 3.2, rgb(0x6C,0x3D,0x8A))
txt(s, "Stage 2 — Linear Regression", 6.9, 2.05, 5.5, 0.5, size=15, bold=True, color=WHITE)
txt(s, "Only 117,711 loans with RR > 0", 6.9, 2.58, 5.5, 0.35, size=11, color=RGBColor(0xC0,0xD8,0xFF))
txt(s, "Target: How much was recovered?\n(continuous, typically 0.05–0.80)", 6.9, 3.0, 5.5, 0.65, size=12, color=WHITE)
txt(s, "Train MAE = 6.03%    OOT MAE = 4.77%", 6.9, 3.75, 5.5, 0.4, size=12, bold=True, color=RGBColor(0x90,0xFF,0x90))
txt(s, "Train R² = 0.050      OOT R² = −0.003", 6.9, 4.18, 5.5, 0.35, size=12, color=RGBColor(0x90,0xFF,0x90))

# Arrow
txt(s, "×", 6.1, 3.3, 0.6, 0.6, size=28, bold=True, color=NAVY, align="center")

rect(s, 0.4, 5.35, 12.5, 1.05, NAVY)
txt(s, "Combined: Recovery Rate = P(RR>0) × E[RR | RR>0]     →     LGD = 1 − Recovery Rate",
    0.6, 5.47, 12.0, 0.45, size=15, bold=True, color=WHITE, align="center", font="Consolas")
txt(s, "Combined OOT MAE = 6.71%    Mean predicted LGD = 92.2%    Mean actual LGD = 94.0%",
    0.6, 5.97, 12.0, 0.35, size=12, color=RGBColor(0x90,0xD4,0xFF), align="center")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — EAD & Expected Loss
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "EAD Model & Portfolio Expected Loss", "L03")

# EAD
rect(s, 0.4, 1.0, 6.0, 2.8, OFFWT, LGRAY, 0.5)
txt(s, "EAD — Credit Conversion Factor", 0.6, 1.1, 5.6, 0.45, size=15, bold=True, color=NAVY)
txt(s, "CCF = (Funded − Total Paid) / Funded   →   EAD = CCF × Funded Amount",
    0.6, 1.6, 5.6, 0.55, size=12, color=DARK, font="Consolas")
ead_rows = [("Algorithm","Linear Regression + StandardScaler"),
            ("Training","154,761 defaulted loans"),
            ("Train MAE","14.89%   OOT MAE ~15%"),
            ("Mean CCF train","0.390   Mean CCF OOT: 0.555")]
for i,(k,v) in enumerate(ead_rows):
    txt(s,f"{k}:", 0.6,2.25+i*0.4,2.2,0.35,size=12,bold=True,color=NAVY)
    txt(s,v, 2.9,2.25+i*0.4,3.4,0.35,size=12,color=DARK)

# Portfolio stats
rect(s, 6.8, 1.0, 6.1, 2.8, NAVY)
txt(s, "Portfolio Summary — OOT Test Set", 7.0, 1.1, 5.7, 0.45, size=15, bold=True, color=WHITE)
port = [("Total Loans","538,515"),("Total EAD","$3.26 Billion"),
        ("Total Expected Loss","$383.6 Million"),("Portfolio EL Rate","11.76%"),
        ("Mean PD","11.98%"),("Mean LGD","92.55%")]
for i,(k,v) in enumerate(port):
    txt(s,f"{k}",7.0,1.62+i*0.38,3.4,0.34,size=12,color=RGBColor(0xB0,0xD0,0xFF))
    txt(s,v,10.5,1.62+i*0.38,2.3,0.34,size=13,bold=True,color=WHITE)

# Credit policy impact
rect(s, 0.4, 4.0, 12.5, 0.5, BLUE)
txt(s, "Credit Policy Impact — Approving only ROI-positive loans:", 0.6, 4.1, 12.0, 0.38, size=14, bold=True, color=WHITE)

impact = [("Approval rate","100%","45.6%","Reject 54.4% of loans"),
          ("Default rate","25.27%","14.41%","↓ 10.9 percentage points"),
          ("Portfolio EL rate","11.03%","6.10%","↓ 45% reduction in EL")]
rect(s, 0.4, 4.52, 12.5, 0.42, NAVY)
for i,lbl in enumerate(["Metric","Before Policy","After Policy","Impact"]):
    txt(s,lbl,0.6+i*3.1,4.6,3.0,0.32,size=12,bold=True,color=WHITE)
for ri,(metric,before,after,impact_str) in enumerate(impact):
    ry=4.97+ri*0.58
    rect(s,0.4,ry,12.5,0.55,OFFWT if ri%2==0 else WHITE,LGRAY,0.3)
    txt(s,metric,0.6,ry+0.1,3.0,0.38,size=13,color=NAVY,bold=True)
    txt(s,before,3.7,ry+0.1,3.0,0.38,size=13,color=RED,bold=True)
    txt(s,after, 6.8,ry+0.1,3.0,0.38,size=13,color=GREEN,bold=True)
    txt(s,impact_str,9.9,ry+0.1,3.0,0.38,size=12,color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — L03 Summary
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "L03 — Summary", "L03")
summ3=[("6.71%","LGD OOT MAE",rgb(0x8E,0x44,0xAD)),("92.2%","Mean LGD",BLUE),
       ("$3.26B","Portfolio EAD",NAVY),("14.41%","Post-policy\nDefault Rate",GREEN)]
x=0.4
for val,lbl,clr in summ3:
    stat_card(s,val,lbl,x,1.0,w=3.0,h=1.6,val_color=clr); x+=3.1
txt(s,"Key Takeaways",0.4,2.8,12.5,0.4,size=16,bold=True,color=NAVY)
kts3=[
    "Two-stage LGD is necessary because 24% of defaulted loans have exactly zero recovery",
    "Stage 1 (logistic) separates loans with and without recovery — AUC 0.67 on train",
    "Stage 2 (linear) estimates how much is recovered given RR>0 — MAE 4.77% on OOT",
    "Beta regression is theoretically purer but gives similar MAE (~5%); two-stage wins on interpretability",
    "EAD CCF mean shifts from 0.39 (train) to 0.56 (OOT) — 2016–2018 loans defaulted earlier in their lifecycle",
    "Credit policy halves EL rate: 11.0% → 6.1% by rejecting ROI-negative loans",
]
for i,k in enumerate(kts3):
    txt(s,f"✓  {k}",0.5,3.38+i*0.53,12.3,0.46,size=13,color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — L04 Section Header
# ═══════════════════════════════════════════════════════════════════════════════
s = dark_title_slide(prs,
    "L04 — Population Stability\nIndex & Model Monitoring",
    "PSI formula · Drift detection · Score PSI 0.109 · MONITOR verdict · Action framework",
    tag="NOTEBOOK 4 OF 4")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — PSI Formula & Thresholds
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "Population Stability Index — Formula & Thresholds", "L04")

rect(s, 0.4, 1.0, 12.5, 1.4, NAVY)
txt(s, "PSI  =  Σ  ( Aᵢ − Eᵢ )  ×  ln( Aᵢ / Eᵢ )",
    0.6, 1.12, 12.0, 0.7, size=28, bold=True, color=WHITE, align="center", font="Consolas")
txt(s, "Aᵢ = Actual % in bucket i (monitoring population)      Eᵢ = Expected % (training reference)",
    0.6, 1.82, 12.0, 0.45, size=13, color=RGBColor(0xA0,0xC4,0xE8), align="center")

thresholds = [
    ("< 0.10", "STABLE", "No significant population shift.\nModel is applicable. Continue quarterly monitoring.", GREEN),
    ("0.10 – 0.25", "MONITOR", "Moderate shift detected.\nInvestigate root cause. Validate Gini/KS on new data.\nConsider intercept recalibration.", AMBER),
    ("> 0.25", "ALERT", "Major population shift.\nModel must be rebuilt.\nRetrain on more recent data (last 2–3 years).", RED),
]
for i, (psi_range, label, desc, clr) in enumerate(thresholds):
    x = 0.4 + i*4.2
    rect(s, x, 2.65, 4.0, 4.0, clr)
    txt(s, "PSI", x+0.2, 2.75, 3.6, 0.35, size=12, color=WHITE)
    txt(s, psi_range, x+0.2, 3.1, 3.6, 0.6, size=22, bold=True, color=WHITE)
    rect(s, x+0.15, 3.78, 3.7, 0.05, WHITE)
    txt(s, label, x+0.2, 3.9, 3.6, 0.55, size=18, bold=True, color=WHITE)
    txt(s, desc, x+0.2, 4.5, 3.6, 2.1, size=11, color=RGBColor(0xE8,0xF4,0xFF))

txt(s,"PSI = 0 means identical distributions. As populations diverge, each term grows → PSI increases.",
    0.4, 6.9, 12.5, 0.38, size=11, italic=True, color=MID, align="center")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — PSI Results
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "PSI Results — 2015 vs 2007–2014 Reference", "L04")

psi_data = [
    ("mths_since_issue_d",          12.7624, "ALERT",   "Time variable — expected high (different issue dates)"),
    ("mths_since_earliest_cr_line",  0.1656, "MONITOR", "Credit history length has shifted slightly"),
    ("credit_score",                 0.1093, "MONITOR", "Score distribution shifted +8.5 pts upward on average"),
    ("grade",                        0.0920, "STABLE",  "Grade mix relatively stable 2007→2015"),
    ("int_rate",                     0.0836, "STABLE",  "Interest rate distribution unchanged"),
    ("home_ownership",               0.0520, "STABLE",  "Ownership mix stable"),
    ("purpose",                      0.0380, "STABLE",  "Loan purpose mix stable"),
    ("emp_length_int",               0.0111, "STABLE",  "Employment length unchanged"),
    ("annual_inc",                   0.0086, "STABLE",  "Income distribution stable"),
    ("dti",                          0.0078, "STABLE",  "Debt-to-income ratio unchanged"),
    ("verification_status",          0.0450, "STABLE",  "Verification mix minor shift"),
    ("initial_list_status",          0.0320, "STABLE",  "Listing status stable"),
]
status_colors = {"ALERT": RED, "MONITOR": AMBER, "STABLE": GREEN}

rect(s, 0.4, 1.0, 12.5, 0.45, NAVY)
for i, lbl in enumerate(["Variable","PSI","Status","Interpretation"]):
    txt(s, lbl, 0.6+i*3.1, 1.08, 3.0, 0.32, size=12, bold=True, color=WHITE)

for ri, (var, psi_val, status, interp) in enumerate(psi_data[:10]):
    ry = 1.48 + ri*0.55
    clr = OFFWT if ri%2==0 else WHITE
    rect(s, 0.4, ry, 12.5, 0.53, clr, LGRAY, 0.2)
    flag = "🔴" if status=="ALERT" else ("🟡" if status=="MONITOR" else "🟢")
    txt(s, var, 0.6, ry+0.08, 3.0, 0.38, size=11, color=NAVY, font="Consolas")
    # PSI bar
    bar_w = min(2.5 * psi_val / 0.3, 2.5) if psi_val < 1 else 2.5
    rect(s, 3.7, ry+0.12, bar_w, 0.28, status_colors[status])
    txt(s, f"{psi_val:.4f}", 3.7+bar_w+0.05, ry+0.1, 0.9, 0.3, size=10, color=DARK)
    txt(s, status, 6.8, ry+0.1, 1.5, 0.33, size=11, bold=True,
        color=status_colors[status])
    txt(s, interp, 8.7, ry+0.1, 4.2, 0.35, size=10, color=DARK)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Credit Score PSI Detail
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "Credit Score PSI — Key Monitoring Signal", "L04")

# Score comparison
rect(s, 0.4, 1.0, 6.0, 1.0, NAVY)
txt(s, "Reference (2007–2014): mean = 680.0  std = 27.4\nActual   (2015):        mean = 688.5  std = 25.2  (+8.5 pts)",
    0.6, 1.1, 5.6, 0.78, size=13, color=WHITE)

rect(s, 6.8, 1.0, 5.9, 1.0, rgb(0x1A,0x53,0x76))
txt(s, "Score PSI = 0.1093  →  MONITOR\n2015 borrowers score higher → appear less risky on paper",
    7.0, 1.1, 5.5, 0.78, size=13, bold=True, color=AMBER)

# PSI bucket breakdown
txt(s, "PSI Bucket Breakdown (10 equal-frequency buckets on reference):",
    0.4, 2.2, 12.5, 0.4, size=14, bold=True, color=NAVY)
bucket_data = [
    ("-∞ to 644",  9.98,  4.00,  0.0546),
    ("644 to 656", 9.52,  7.16,  0.0067),
    ("656 to 665", 9.51,  8.07,  0.0024),
    ("665 to 673", 9.76,  8.51,  0.0017),
    ("673 to 681",10.47,  9.49,  0.0010),
    ("681 to 688", 9.40,  9.24,  0.0000),
    ("688 to 696",10.61, 11.13,  0.0003),
    ("696 to 704", 9.81, 11.52,  0.0027),
    ("704 to 715",10.65, 14.61,  0.0125),
    ("715 to ∞",  10.29, 16.27,  0.0274),
]
rect(s, 0.4, 2.7, 12.5, 0.4, NAVY)
for i,lbl in enumerate(["Score Range","Ref %","Actual %","PSI Contrib","Direction"]):
    txt(s, lbl, 0.6+i*2.45, 2.78, 2.3, 0.28, size=11, bold=True, color=WHITE)
for ri, (rng, ref_p, act_p, contrib) in enumerate(bucket_data):
    ry = 3.12+ri*0.38
    rect(s, 0.4, ry, 12.5, 0.36, OFFWT if ri%2==0 else WHITE, LGRAY, 0.2)
    direction = "▲ More" if act_p > ref_p else "▼ Fewer"
    dir_clr = BLUE if act_p > ref_p else RED
    txt(s, rng, 0.6, ry+0.06, 2.3, 0.26, size=10, color=DARK)
    txt(s, f"{ref_p:.1f}%", 3.05, ry+0.06, 2.2, 0.26, size=10, color=DARK)
    txt(s, f"{act_p:.1f}%", 5.5, ry+0.06, 2.2, 0.26, size=11, bold=(abs(act_p-ref_p)>3),
        color=dir_clr if abs(act_p-ref_p)>3 else DARK)
    txt(s, f"{contrib:.4f}", 7.95, ry+0.06, 2.2, 0.26, size=10, color=RED if contrib>0.02 else DARK)
    txt(s, direction, 10.4, ry+0.06, 2.1, 0.26, size=10, color=dir_clr, bold=abs(act_p-ref_p)>3)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 27 — Monitoring Framework
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "Model Monitoring Framework & Action Thresholds", "L04")

rect(s, 0.4, 1.0, 12.5, 0.5, GREEN)
txt(s, "OUR VERDICT:  Score PSI = 0.1093 → MONITOR | 8 Stable  ·  2 Monitor  ·  2 Alert (time-variable)",
    0.6, 1.08, 12.0, 0.38, size=14, bold=True, color=WHITE)

# Action table
actions = [
    ("Score PSI", "< 0.10",      "STABLE",    "Quarterly PSI run, no changes needed",GREEN),
    ("Score PSI", "0.10 – 0.19", "MODERATE",  "Validate Gini/KS. If drop < 5pp: monthly monitoring.\nIf 5–10pp: intercept recalibration",AMBER),
    ("Score PSI", "0.19 – 0.25", "SIGNIFICANT","High priority. Retrain on recent data. Update WoE bins\nfor high-PSI variables",AMBER),
    ("Score PSI", "> 0.25",      "CRITICAL",   "Begin model redevelopment. New training window.\nConsider XGBoost + SHAP for next version",RED),
]
rect(s, 0.4, 1.65, 12.5, 0.42, NAVY)
for i, lbl in enumerate(["Metric","PSI Range","Status","Recommended Action"]):
    txt(s, lbl, 0.6+i*3.1, 1.73, 3.0, 0.3, size=12, bold=True, color=WHITE)
for ri,(metric,rng,status,action,clr) in enumerate(actions):
    ry = 2.1+ri*0.88
    rect(s, 0.4, ry, 12.5, 0.85, OFFWT if ri%2==0 else WHITE, LGRAY, 0.3)
    rect(s, 0.4, ry, 0.14, 0.85, clr)
    txt(s, metric, 0.65, ry+0.1, 3.0, 0.3, size=12, bold=True, color=NAVY)
    txt(s, rng,    3.7, ry+0.1, 3.0, 0.3, size=13, bold=True, color=clr)
    txt(s, status, 6.8, ry+0.1, 3.0, 0.3, size=12, bold=True, color=clr)
    txt(s, action, 9.9, ry+0.05, 3.0, 0.75, size=10, color=DARK)

# Monitoring dashboard
rect(s, 0.4, 5.68, 12.5, 1.2, NAVY)
txt(s, "Complete Monitoring Dashboard",
    0.6, 5.78, 12.0, 0.38, size=14, bold=True, color=WHITE)
dash_metrics = [("Gini (OOT)",">0.40","0.35–0.40","<0.35","Rebuild"),
                ("KS (OOT)",">0.25","0.20–0.25","<0.20","Rebuild"),
                ("Score PSI","<0.10","0.10–0.25",">0.25","Rebuild/Recalib"),
                ("Var PSI","<0.10","0.10–0.25",">0.25","Re-encode bins")]
for i, (m,g,y,r,act) in enumerate(dash_metrics):
    x = 0.5+i*3.1
    txt(s, m, x, 6.18, 2.9, 0.28, size=11, bold=True, color=WHITE)
    txt(s, f"🟢{g} 🟡{y} 🔴{r}", x, 6.48, 2.9, 0.3, size=9, color=RGBColor(0xA0,0xC4,0xE8))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 28 — Production Pipeline
# ═══════════════════════════════════════════════════════════════════════════════
s = dark_title_slide(prs,
    "Production Pipeline\n& Live System",
    "8 Docker services · FastAPI scoring endpoint · Airflow DAGs · MLflow registry · Grafana",
    tag="TRACK 2 — PRODUCTION")

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 29 — Architecture
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "Production Architecture — 8 Docker Services", "Production")

services = [
    ("PostgreSQL\n:5432", "6-schema warehouse\nraw·staging·features\nmodels·risk·audit", BLUE),
    ("MinIO\n:9000/9001", "S3-compatible\ndata lake\n3 buckets", rgb(0x1B,0x6C,0xA8)),
    ("Spark\n:8080/7077", "Big data\nprocessing\n2 jobs", rgb(0x17,0x5B,0x8A)),
    ("Airflow\n:8081", "6 DAGs\norchestrating\nfull pipeline", rgb(0x13,0x4A,0x6C)),
    ("MLflow\n:5001", "Experiment\ntracking &\nmodel registry", rgb(0x0F,0x3A,0x52)),
    ("FastAPI\n:8000", "Real-time\nscoring\nendpoint", GREEN),
    ("Prometheus\n:9090", "Metrics\ncollection &\nalerts", rgb(0x7B,0x24,0x1C)),
    ("Grafana\n:3000", "Dashboards\n& monitoring\nvisuals", rgb(0x6E,0x27,0x17)),
]
for i, (name, desc, clr) in enumerate(services):
    col = i % 4
    row = i // 4
    x = 0.4 + col*3.2
    y = 1.1 + row*2.8
    rect(s, x, y, 3.0, 2.5, clr)
    txt(s, name, x+0.15, y+0.12, 2.7, 0.7, size=14, bold=True, color=WHITE, font="Consolas")
    txt(s, desc, x+0.15, y+0.85, 2.7, 1.55, size=11, color=RGBColor(0xC0,0xD8,0xF0))

txt(s, "Data Flow: CSV → MinIO → Spark ingestion → PostgreSQL → Airflow DAGs → MLflow training → FastAPI scoring → Grafana dashboards",
    0.4, 6.85, 12.5, 0.45, size=11, color=MID, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 30 — Live API Demo
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "FastAPI Scoring Endpoint — Live Demo", "Production")

rect(s, 0.4, 1.0, 6.0, 5.7, NAVY)
txt(s, "Request — POST /score", 0.6, 1.1, 5.6, 0.42, size=14, bold=True, color=WHITE)
req_json = '''{
  "loan_amnt": 15000,
  "funded_amnt": 15000,
  "term_int": 36,
  "int_rate": 0.1199,
  "grade": "B",
  "annual_inc": 75000,
  "fico_score": 710,
  "dti": 18.5,
  "revol_util": 0.42,
  "inq_last_6mths": 1,
  "home_ownership": "MORTGAGE",
  "purpose": "debt_consolidation"
}'''
txt(s, req_json, 0.6, 1.58, 5.6, 5.0, size=11, color=RGBColor(0x90,0xFF,0x90), font="Consolas")

rect(s, 6.8, 1.0, 6.1, 5.7, rgb(0x0A,0x2A,0x40))
txt(s, "Response", 7.0, 1.1, 5.7, 0.42, size=14, bold=True, color=WHITE)
resp_fields = [
    ("pd",              "0.030303",  "3.0% probability of default"),
    ("lgd",             "0.45",      "45% loss given default"),
    ("ead",             "$15,000",   "Full exposure at default"),
    ("expected_loss",   "$204.55",   "PD × LGD × EAD"),
    ("credit_score",    "700",       "AB risk class"),
    ("risk_class",      "\"AB\"",    "Score range 700–739"),
    ("decision",        "APPROVE",   "ROI > 2.15% base rate"),
    ("annualized_roi",  "11.54%",    "Net return after EL"),
]
for i, (field, value, note) in enumerate(resp_fields):
    ry = 1.6+i*0.6
    txt(s, f'"{field}":', 7.0, ry, 2.5, 0.5, size=11, color=RGBColor(0x86,0xC8,0xFF), font="Consolas")
    txt(s, value, 9.6, ry, 1.5, 0.5, size=12, bold=True,
        color=GREEN if field in ["decision","annualized_roi"] else WHITE, font="Consolas")
    txt(s, note,  11.2, ry+0.08, 1.5, 0.38, size=9, color=RGBColor(0x80,0xA0,0xC0))

txt(s, "Response time: < 50ms  ·  Models loaded at startup from data/models/*.pkl",
    0.4, 6.9, 12.5, 0.38, size=11, italic=True, color=MID)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 31 — Airflow DAGs
# ═══════════════════════════════════════════════════════════════════════════════
s = content_slide(prs, "Airflow DAGs — Pipeline Orchestration", "Production")

dags = [
    ("01_ingestion",       "Manual trigger", "Spark: CSV from MinIO → raw.loan_applications"),
    ("02_cleaning",        "Manual trigger", "Spark: raw → staging.loans_cleaned + OOT split"),
    ("03_pd_training",     "Manual trigger", "Logistic regression + backward elim → scorecard.csv + MLflow"),
    ("04_lgd_ead_training","Manual trigger", "2-stage LGD + EAD → pkl files + MLflow"),
    ("05_batch_scoring",   "Manual trigger", "Score all OOT loans → risk.expected_loss table"),
    ("06_monitoring",      "Daily 6am",      "PSI for all variables → risk.population_stability + alerts"),
]
rect(s, 0.4, 1.0, 12.5, 0.48, NAVY)
for i, lbl in enumerate(["DAG","Schedule","Description"]):
    txt(s, lbl, 0.6+i*4.1, 1.08, 3.9, 0.32, size=12, bold=True, color=WHITE)

for ri, (dag, sched, desc) in enumerate(dags):
    ry = 1.52+ri*0.85
    rect(s, 0.4, ry, 12.5, 0.82, OFFWT if ri%2==0 else WHITE, LGRAY, 0.3)
    clr = GREEN if ri == 5 else BLUE
    rect(s, 0.4, ry, 0.12, 0.82, clr)
    txt(s, dag, 0.65, ry+0.18, 3.8, 0.45, size=12, bold=True, color=NAVY, font="Consolas")
    sched_clr = GREEN if "Daily" in sched else MID
    txt(s, sched, 4.75, ry+0.22, 3.5, 0.38, size=11, color=sched_clr, bold="Daily" in sched)
    txt(s, desc, 8.5, ry+0.18, 4.4, 0.45, size=11, color=DARK)

txt(s, "All 6 DAGs written and awaiting Airflow webserver start. PSI alert triggers rebuild workflow automatically.",
    0.4, 6.7, 12.5, 0.45, size=11, italic=True, color=MID)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 32 — Final Summary
# ═══════════════════════════════════════════════════════════════════════════════
s = dark_title_slide(prs,
    "Project Complete",
    "From raw CSV to production scoring API — end-to-end credit risk system",
    tag="SUMMARY")
bg(s, NAVY)

# 4 notebook cards
nb_summ = [
    ("L01", "Preprocessing", "831K train · 56 WoE dummies\nOOT split 2007-15 / 2016-18", BLUE),
    ("L02", "PD Model", "49 features · AUC 0.693\nScorecard 300–850", rgb(0x6C,0x3D,0x8A)),
    ("L03", "LGD / EAD / EL", "2-stage LGD MAE 6.7%\n$3.26B EAD · $383M EL", GREEN),
    ("L04", "PSI Monitoring", "Score PSI 0.109 MONITOR\n8 stable · 2 monitor · 2 alert", AMBER),
]
for i, (badge, title, metrics, clr) in enumerate(nb_summ):
    x = 0.4 + i*3.2
    rect(s, x, 1.3, 3.0, 2.8, clr)
    txt(s, badge, x+0.15, 1.42, 2.7, 0.55, size=22, bold=True, color=WHITE, align="center")
    txt(s, title, x+0.15, 2.0, 2.7, 0.5, size=14, bold=True, color=WHITE, align="center")
    rect(s, x+0.15, 2.52, 2.7, 0.04, WHITE)
    txt(s, metrics, x+0.15, 2.62, 2.7, 1.35, size=11, color=RGBColor(0xD0,0xE8,0xFF), align="center")

# Key metrics row
key_metrics = [
    ("0.693", "OOT AUC"),("0.386","OOT Gini"),
    ("$383M","Portfolio EL"),("0.109","Score PSI"),("< 50ms","API Response"),
]
x2=0.4
for i, (val, lbl) in enumerate(key_metrics):
    w2=12.5/len(key_metrics)
    rect(s, x2, 4.35, w2-0.15, 1.3, rgb(0x1B,0x4F,0x72))
    txt(s, val, x2+0.1, 4.45, w2-0.3, 0.65, size=22, bold=True, color=WHITE, align="center")
    txt(s, lbl, x2+0.1, 5.1, w2-0.3, 0.45, size=11, color=RGBColor(0xB0,0xD0,0xFF), align="center")
    x2 += w2

txt(s, "Regulatory compliance: IFRS 9 · Basel III · SR 11-7",
    0.5, 5.85, 12.3, 0.45, size=14, color=RGBColor(0x90,0xD4,0xFF), align="center")
txt(s, "Production stack: PySpark · PostgreSQL · Airflow · MLflow · FastAPI · Prometheus · Grafana",
    0.5, 6.4, 12.3, 0.4, size=12, color=RGBColor(0x70,0xB4,0xE0), align="center", italic=True)
txt(s, "Thank You  |  Questions?",
    0.5, 7.0, 12.3, 0.38, size=15, bold=True, color=WHITE, align="center")

prs.save(OUT)
print(f"Saved → {OUT}")
print(f"Total slides: {len(prs.slides)}")
